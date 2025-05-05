from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import os
from PIL import Image
import re
from difflib import SequenceMatcher

class PPTGenerator:
    def __init__(self):
        self.used_images = set()
    
    def generate(self, slides_content, images, output_path):
        print(f"Number of images available: {len(images)}")
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_shape = title_slide.shapes.title
        subtitle_shape = title_slide.placeholders[1] if len(title_slide.placeholders) > 1 else None

        if slides_content:
            title_shape.text = slides_content[0]["title"]
            if subtitle_shape:
                subtitle_shape.text = "Document Summary"

        # Find a suitable image for the title slide
        title_image = self._find_title_image(images)
        if title_image:
            self._add_image_to_slide(title_slide, title_image, is_title_slide=True)

        # Content Slides
        for idx, slide_content in enumerate(slides_content[1:], 1):
            print(f"\nProcessing slide {idx}: {slide_content['title']}")
            
            # Safely get image_hint and bullets
            image_hint = slide_content.get("image_hint", "")
            if image_hint is not None:
                image_hint = image_hint.lower()
            else:
                image_hint = ""
                
            bullets = slide_content.get("bullets", [])
            if bullets is None:
                bullets = []
            
            matching_image = self._find_matching_image(image_hint, images, slide_content["title"], bullets)
            
            # Use different layouts based on whether we have an image
            if matching_image:
                # Two Content layout (for text + image)
                slide_layout = prs.slide_layouts[3]  # Usually layout 3 is for two content
            else:
                # Title and Content layout (just text)
                slide_layout = prs.slide_layouts[1]
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = slide_content["title"]
                
                # Format title text
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(40)
                        run.font.bold = True
            
            # Add content based on layout
            if matching_image:
                # Two content placeholders - left for text, right for image
                placeholders = [shape for shape in slide.placeholders 
                               if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
                               and shape.placeholder_format.idx != 0]  # Skip title placeholder
                
                if len(placeholders) >= 2:
                    # Left placeholder for text
                    text_placeholder = placeholders[0]
                    text_frame = text_placeholder.text_frame
                    text_frame.clear()
                    
                    for bullet in bullets:
                        p = text_frame.add_paragraph()
                        p.text = bullet.strip()
                        p.level = 0
                        # Format bullet text
                        for run in p.runs:
                            run.font.size = Pt(24)
                    
                    # Right placeholder for image
                    self._add_image_to_placeholder(slide, placeholders[1], matching_image)
                    self.used_images.add(matching_image["path"])
                    print(f"Added image: {matching_image['path']}")
                else:
                    # Fallback if layout doesn't have expected placeholders
                    self._add_content_with_custom_image(slide, slide_content, matching_image)
            else:
                # Just add text content
                content_shape = None
                for shape in slide.placeholders:
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx != 0:
                        content_shape = shape
                        break
                
                if content_shape:
                    text_frame = content_shape.text_frame
                    text_frame.clear()
                    
                    for bullet in bullets:
                        p = text_frame.add_paragraph()
                        p.text = bullet.strip()
                        p.level = 0
                        # Format bullet text
                        for run in p.runs:
                            run.font.size = Pt(24)

        # Save the presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        print(f"\nPresentation saved to: {output_path}")
        print(f"Used {len(self.used_images)} images out of {len(images)} available")

    def _add_image_to_placeholder(self, slide, placeholder, image_info):
        """Add image to a placeholder, maintaining aspect ratio"""
        img_path = image_info["path"]
        if not os.path.exists(img_path):
            print(f"Image file not found: {img_path}")
            return
        
        try:
            # Get placeholder dimensions
            ph_width = placeholder.width
            ph_height = placeholder.height
            
            # Calculate image dimensions to maintain aspect ratio
            with Image.open(img_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
                
                if (ph_width / ph_height) > aspect_ratio:
                    # Placeholder is wider than image proportionally
                    new_height = ph_height
                    new_width = int(new_height * aspect_ratio)
                else:
                    # Placeholder is taller than image proportionally
                    new_width = ph_width
                    new_height = int(new_width / aspect_ratio)
            
            # Center the image in the placeholder
            left = int(placeholder.left + (ph_width - new_width) / 2)
            top = int(placeholder.top + (ph_height - new_height) / 2)
            
            # Add image to slide (not to the placeholder directly)
            slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)
            
        except Exception as e:
            print(f"Failed to add image to placeholder: {str(e)}")

    def _add_image_to_slide(self, slide, image_info, is_title_slide=False):
        """Add image to slide with custom positioning"""
        img_path = image_info["path"]
        if not os.path.exists(img_path):
            print(f"Image file not found: {img_path}")
            return
        
        try:
            # Get slide dimensions
            slide_width = slide.slide_layout.width
            slide_height = slide.slide_layout.height
            
            if is_title_slide:
                # For title slide, place image at bottom right
                max_width = Inches(5)  # Maximum width for the image
                
                with Image.open(img_path) as img:
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                    
                    new_width = max_width
                    new_height = int(new_width / aspect_ratio)
                    
                # Position at bottom right with some margin
                left = int(slide_width - new_width - Inches(0.5))
                top = int(slide_height - new_height - Inches(0.5))
                
            else:
                # For content slides, place on right half
                max_width = int(slide_width / 2 - Inches(0.5))
                
                with Image.open(img_path) as img:
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                    
                    new_width = max_width
                    new_height = int(new_width / aspect_ratio)
                    
                    if new_height > slide_height - Inches(2):
                        # Too tall, scale by height instead
                        new_height = int(slide_height - Inches(2))
                        new_width = int(new_height * aspect_ratio)
                
                # Position on right half
                left = int(slide_width / 2 + Inches(0.25))
                top = int(Inches(1.5))
            
            # Add image to slide
            slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)
            self.used_images.add(img_path)
            print(f"Added image: {img_path}")
            
        except Exception as e:
            print(f"Failed to add image to slide: {str(e)}")

    def _add_content_with_custom_image(self, slide, slide_content, image_info):
        """Add text content on left side and image on right side manually"""
        # Add bullets as text box on left side
        left = int(Inches(0.5))
        top = int(Inches(1.5))
        width = int(Inches(6))
        height = int(Inches(5))
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txBox.text_frame
        
        # Safe handling of bullets
        bullets = slide_content.get("bullets", [])
        if bullets is None:
            bullets = []
            
        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = "• " + bullet.strip()  # Add bullet character manually
            # Format bullet text
            for run in p.runs:
                run.font.size = Pt(24)
            
            # Add spacing between bullets
            p.space_after = Pt(12)
        
        # Add image on right side
        self._add_image_to_slide(slide, image_info)

    def _find_title_image(self, images):
        """Find a suitable image for the title slide"""
        if not images:
            return None
            
        # Look for logo or title-related images first
        logo_keywords = ['logo', 'title', 'cover', 'header', 'main']
        for image in images:
            context = image.get("context", "")
            if context is None:
                context = ""
            else:
                context = context.lower()
                
            if any(keyword in context for keyword in logo_keywords):
                return image
                
        # If no specific title image found, use the largest image
        largest_image = None
        max_size = 0
        
        for image in images:
            if "size" in image:
                width, height = image["size"]
                size = width * height
                if size > max_size:
                    max_size = size
                    largest_image = image
        
        # If we still don't have an image, just use the first one
        if not largest_image and images:
            largest_image = images[0]
            
        return largest_image

    def _find_matching_image(self, image_hint, images, slide_title, bullets):
        """Find the most contextually relevant image for the slide"""
        if not images:
            return None

        # Ensure we have string values
        if slide_title is None:
            slide_title = ""
            
        # Create a combined context from title and bullets for better matching
        if bullets:
            bullet_text = " ".join(str(b) for b in bullets if b is not None)
        else:
            bullet_text = ""
            
        slide_context = slide_title.lower() + " " + bullet_text.lower()
        
        # Extract figure/table numbers from hints
        fig_numbers = []
        if image_hint:
            fig_match = re.search(r'(figure|fig\.?|table)\s*(\d+)', image_hint)
            if fig_match:
                fig_numbers.append(fig_match.group(2))
        
        # Extract figure/table numbers from slide content
        for text in [slide_title] + (bullets if bullets else []):
            if text is None:
                continue
            matches = re.finditer(r'(figure|fig\.?|table)\s*(\d+)', text.lower())
            for match in matches:
                fig_numbers.append(match.group(2))
        
        best_match = None
        best_score = 0
        
        def text_similarity(text1, text2):
            """Calculate the similarity between two text strings"""
            if not text1 or not text2:
                return 0
            return SequenceMatcher(None, text1, text2).ratio()
        
        # First try to match by figure/table number
        if fig_numbers:
            for image in images:
                if image["path"] in self.used_images:
                    continue
                    
                context = image.get("context", "")
                if context is None:
                    context = ""
                else:
                    context = context.lower()
                    
                for num in fig_numbers:
                    if re.search(rf'(figure|fig\.?|table)\s*{num}\b', context):
                        return image
        
        # If no exact figure match, try semantic matching
        for image in images:
            if image["path"] in self.used_images:
                continue
                
            context = image.get("context", "")
            if context is None:
                context = ""
            elif context:
                context = context.lower()
            else:
                # If no context, use filename as context
                context = os.path.basename(image["path"]).lower()
            
            # Calculate similarity score
            similarity = text_similarity(slide_context, context)
            
            # Give bonus points for keywords
            keywords = ['chart', 'graph', 'plot', 'diagram', 'screenshot', 'illustration', 
                       'figure', 'table', 'image', 'photo', 'picture']
            
            # Check if any keywords appear in both slide context and image context
            for keyword in keywords:
                if keyword in slide_context and keyword in context:
                    similarity += 0.2
            
            if similarity > best_score:
                best_score = similarity
                best_match = image
        
        # Only use images with a reasonable match score
        if best_score > 0.3:
            return best_match
            
        # If we couldn't find a good match but have unused images, pick the first unused one
        if not best_match:
            for image in images:
                if image["path"] not in self.used_images:
                    return image
                    
        return best_match